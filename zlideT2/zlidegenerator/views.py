import os
import json
from random import randint
from django.http import HttpResponse
from rest_framework.views import APIView
from rest_framework.permissions import AllowAny
from rest_framework.response import Response
from rest_framework.exceptions import NotFound
from rest_framework.generics import GenericAPIView
from rest_framework import status
from drf_spectacular.types import OpenApiTypes
from drf_spectacular.utils import extend_schema
from pptx import Presentation
import nltk
from nltk import word_tokenize
from nltk import pos_tag
from .models import PresentationData
from .serializers import PresentationDataSerializer
from openai import OpenAI


nltk.download('punkt')
nltk.download('averaged_perceptron_tagger')

# Setting the NLTK data path
nltk.data.path.append(os.path.join(os.path.dirname(__file__), 'nltk_data'))

OpenAI.api_key = os.environ.get("OPENAI_API_KEY")
client = OpenAI()

class GenerateZlideView(APIView):
    permission_classes = [AllowAny]

    def _extract_first_noun_determiner(self, text):
        tokens = word_tokenize(text)
        tagged = pos_tag(tokens)

        noun = determiner = None
        for word, tag in tagged:
            if not noun and tag.startswith('NN'): # Noun 
                noun = word.capitalize() 
            if not determiner and tag.startswith('DT'): # Determiner
                determiner = word.capitalize()
            if noun and determiner:
                break
        return noun, determiner

    def _generate_presentation(self, user_input, request):
        user_input = request.data.get('text')

        if len(user_input) >= 50:
            # Splitting the user input into slides based on "Slide" keyword
            slides = user_input.split("Slide")
            zlide = []
            for idx, slide in enumerate(slides):
                if slide.strip(): # Ignoring strings
                    try:
                        title, content = slide.split(":", 1)
                        title = title.strip()
                        content = content.strip()
                        # Extracting first noun and determiner from the content
                        noun, determiner = self._extract_first_noun_determiner(content)
                        if noun and determiner:
                            title = f"{determiner} {noun}"
                        elif noun:
                            title = noun
                        elif determiner:
                            title = determiner
                        # Creating a JSON object for the slide
                        zlide_json = {
                            "slide": idx,
                            "title": title,
                            "content": content
                        }
                        # Append JSON object to the list
                        zlide.append(zlide_json)
                    except ValueError:
                        continue # Handle any case where the split fails
            zlide_json_output = json.dumps(zlide, indent=4)
            return json.loads(zlide_json_output)
            # Return the list of JSON objects
        else:                    
            prompt = f"Generate a 5 slide content for a Powerpoint presentation with slides, titles and content and convert them into a JSON array with each item having a slide, title, content about: {user_input}"
            completion = client.chat.completions.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": prompt}])
            response = completion.choices[0].message.content
            return json.loads(response)



    @extend_schema(
        operation_id="Generate Zlide Endpoint",
        description="This endpoint does the actual slide generation",
        summary="This endpoint will generate a slide based on user input by making a call to the OpenAI API and then return the data in JSON format",
        request=OpenApiTypes.OBJECT,
        responses={200: PresentationDataSerializer},
    )
    def post(self, request):
        try:
            user_input = request.data.get("user_input")
            input_text = request.data.get("input_text")
            request.data['user_input'] = user_input
            presentation_data = self._generate_presentation(input_text, request)

            # If not isinstance(presentation_data, dict):
            # return Response({"error": "Invalid response from OpenAI"}, status=status.HTTP_400_BAD_REQUEST)

            # presentation_title = request.data.get("title", "Presentation Title")
        
            # serialized_data = json.dumps(presentation_data)
            # presentation_data = PresentationData.objects.create(title=presentation_title, json_data=serialized_data)
            # serializer = PresentationDataSerializer(presentation_data)
            # return Response({'message': 'Presentation data saved successfully.', 'presentation_id': presentation_data.id, 'slide_data':json.loads(serializer.data["json_data"])}, status=status.HTTP_201_CREATED)
            return Response({'message': 'Presentation created successfully.', 'slide_data':presentation_data}, status=status.HTTP_200_OK)
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class SaveZlideView(APIView):
    permission_classes = [AllowAny]
    queryset = PresentationData.objects.all()
    serializer_class = PresentationDataSerializer

    @extend_schema(
        operation_id="Save Zlide to DB Endpoint",
        description="This endpoint will save a slide based on the title specified by the user",
        summary="This endpoint is used to save a slide to the database",
        request=OpenApiTypes.OBJECT,
        responses={200: PresentationDataSerializer},
    )
    def post(self, request):
        presentation_data = request.data.get("presentation_data")
        title = request.data.get("title")
        if not presentation_data or not title:
            return Response({'error': 'Missing required fields: presentation_data and title'}, status=status.HTTP_400_BAD_REQUEST)
        try:
            # Deserializing the presentation data (assuming it's in JSON format)
            deserialized_data = json.loads(presentation_data)
            serializer = PresentationDataSerializer(data={'title': title, 'json_data': deserialized_data})
            if serializer.is_valid():
                serializer.save()
                return Response({'message': 'Presentation data saved successfully.', 'presentation_id': serializer.data['id'], 'presentation_title': serializer.data['title']}, status=status.HTTP_201_CREATED)
            else:
                return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
        except PresentationData.DoesNotExist:
            return Response({'error': 'No presentation data found with the given title'}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        

class GetZlideView(GenericAPIView):
    permission_classes = [AllowAny]
    queryset = PresentationData.objects.all()
    serializer_class = PresentationDataSerializer
    lookup_field = 'title' # Specifying the field to use for the lookup

    @extend_schema(
        operation_id='Get Zlide Endpoint',
        description='This endpoint retrieves a slide based on the title specified by the user',
        summary='This endpoint will get a slide specified by the user',
        request=OpenApiTypes.OBJECT,
        responses={200: PresentationDataSerializer},
    )
    def get(self, request, title):
        try:
            self.kwargs[self.lookup_field] = title 
            presentation_data = self.get_object()
            serializer = self.get_serializer(presentation_data)
            return Response(serializer.data)
        except NotFound:
            return Response({'error': 'No slide found with the given title'}, status=status.HTTP_404_NOT_FOUND)

class DownloadZlideView(APIView):
    def _deserialize_json_data(self, serializer):
        json_data = json.loads(serializer.data["json_data"])
        if not isinstance(json_data, dict) or 'slides' not in json_data:
            raise ValueError("Invalid JSON data format: Expected dictionary with 'slides' key")
        return json_data
    

    def _create_pptx_presentation(self, slide_data):
        """
        Create a PPTX presentation from the slide data.
        """
        prs = Presentation()
        for slide in slide_data:
            if not isinstance(slide, dict) or 'header' not in slide or 'content' not in slide:
                raise ValueError("Invalide slide data format: Expected dictionary with 'header' and 'content' keys")
            slide_layout = prs.slide_layouts[1]  # Use the Title and Content layout
            new_slide = prs.slides.add_slide(slide_layout)
            if slide.get('header'):
                title = new_slide.shapes.title
                title.text = slide['header']
            if slide.get('header'):
                title = new_slide.shapes.title
                title.text = slide['header']
            if slide.get('content'):
                shapes = new_slide.shapes
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = slide["content"]
                tf.fit_text(font_family="Calibri", max_size=18, bold=True)
            # Output_filename = "output.pptx"
            output_filename = title
            prs.save(output_filename)
            

    @extend_schema( 
        operation_id="Download Zlide Endpoint",
        description="This endpoint downloads the slide",
        summary="This endpoint allows the user to download the slide in PPTX format",
        request=OpenApiTypes.OBJECT,
        responses={200: PresentationDataSerializer},
    )
    def get(self, request):
        try:
            presentation_data = PresentationData.objects.latest('created_at')
            serializer = PresentationDataSerializer(presentation_data)
            json_data = self._deserialize_json_data(serializer)
            slide_data = json_data.get("slides", [])
            if not isinstance(slide_data, list) or not all(isinstance(slide, dict) for slide in slide_data):
                raise ValueError("Invalid slide data format: Expected list of dictionaries")
            output_filename = self._create_pptx_presentation(slide_data)
            with open(output_filename, 'rb') as f:
                response = HttpResponse(f.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                response['Content-Disposition'] = 'attachment; filename="output.pptx"'
                return response
        except PresentationData.DoesNotExist:
            return Response({"error": "No presentation data found, are you sure you\'ve created it?"}, status=status.HTTP_404_NOT_FOUND)
        except ValueError as e:
            return Response({"error": str(e)}, status=status.HTTP_400_BAD_REQUEST)
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
    

class EditZlideView(GenericAPIView):
    permission_classes = [AllowAny]
    queryset = PresentationData.objects.all()
    serializer_class = PresentationDataSerializer
    lookup_field = 'title' # Specifying the field to use for the lookup

    @extend_schema(
        operation_id='Edit Zlide Endpoint',
        description='This endpoint updates the slide data items such as the title, content as specified by the user',
        summary='This endpoint will edit a slide specified by the user',
        request=OpenApiTypes.OBJECT,
        responses={200: PresentationDataSerializer},
    )
    def patch(self, request, title):
        try:
            # self.kwargs[self.lookup_field] = title # Setting the lookup field to the title specified by the user
            # presentation_data = self.get_object() # Using the inherited get_object method

            # # Making sure json_data is a list of dictionaries
            # if isinstance(presentation_data.json_data, list):
            #     json_data = json.loads(presentation_data.json_data)
            # else:
            #     json_data = presentation_data.json_data

            # if isinstance(json_data, str):
            #     json_data = json.loads(json_data)

            # updates = request.data.get('json_data', [])

            # # Apply updates to the json_data field
            # for update in updates:
            #     # slide_number = str(update.get('slide'))  # Convert to string if needed
            #     slide_number = update.get('slide')
            #     updated = False
            #     for slide in json_data:
            #         if slide[randint(0, 6)] == slide_number:
            #             slide.update(update)
            #             updated = True
            #             break
            #     if not updated:
            #         return Response({'error': f'Slide {slide_number} not found for update'}, status=status.HTTP_404_NOT_FOUND)

            
            # presentation_data.json_data = json.dumps(json_data) # Convert back to JSON string if necessary
            # presentation_data.save()

            # serializer = self.get_serializer(presentation_data)
            # return Response(serializer.data)

            # 2nd Alternative approach
            self.kwargs[self.lookup_field] = title
            presentation_data = self.get_object()
            serializer = self.get_serializer(presentation_data, data=request.data, partial=True)

            if serializer.is_valid():
                serializer.save()
                return Response({'message': 'Slide data updated successfully.', 'updated_slide': serializer.data}, status=status.HTTP_200_OK)
            else:
                return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

            # 3rd Alternative approach
            # self.kwargs[self.lookup_field] = title
            # presentation_data = self.get_object()
            # deserialized_data = json.loads(presentation_data)
            # # serializer = self.get_serializer(presentation_data)
            # serializer = self.get_serializer(data={'title': title, 'json_data': deserialized_data})

            # # Access and update specific slide data based on request payload
            # update_data = request.data.get('json_data', [])
            # for update_item in update_data:
            #     slide_number = update_item.get('slide')
            #     if slide_number is None:
            #         return Response({'error': 'Missing slide number in update data'}, status=status.HTTP_400_BAD_REQUEST)

            #     # Find the corresponding slide data in the presentation
            #     json_data = next((item for item in presentation_data.json_data if item['slide'] == slide_number), None)
            #     if not json_data:
            #         return Response({'error': f'Slide with number {slide_number} not found'}, status=status.HTTP_404_NOT_FOUND)

            #     # Update the serializer data with the modified slide data
            #     serializer.data[slide_number - 1] = update_item

            #     # Update specific fields (title/content) based on provided data
            #     for field, value in update_item.items():
            #         if field in ['title', 'content']:
            #             json_data[field] = value

            # if serializer.is_valid():
            #     serializer.save()
            #     return Response(serializer.data, status=status.HTTP_200_OK)
            # else:
            #     return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

            return Response(serializer.data)
        except NotFound:
            return Response({'error': f'{title} not found'}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


class DeleteZlideView(APIView):
    permission_classes = [AllowAny]
    queryset = PresentationData.objects.all()
    serializer_class = PresentationDataSerializer

    @extend_schema(
        operation_id='Delete Zlide Endpoint',
        description='This endpoint is used to delete a slide from the database',
        summary='This endpoint will delete a slide specified by the user',
        request=OpenApiTypes.OBJECT,
        responses={200: PresentationDataSerializer},
    )
    def post(self, request):
        try:
            title = request.data.get('title')
            if not title:
                return Response({'error': 'Title is required'}, status=status.HTTP_400_BAD_REQUEST)

            presentation_data = PresentationData.objects.get(title=title)
            presentation_data.delete()
            return Response({'message': f'{title} deleted successfully.'}, status=status.HTTP_200_OK)
        except PresentationData.DoesNotExist:
            return Response({'error': 'No presentation found with the given title'}, status=status.HTTP_404_NOT_FOUND)
        except Exception as e:
            return Response({'error': str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)